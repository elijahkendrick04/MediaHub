"""Document engine (roadmap 1.15) — build 3: PPTX / DOCX exports."""

from __future__ import annotations

from mediahub.documents import export
from mediahub.documents import models as m
from mediahub.documents.models import DocumentSpec, Section


def _doc():
    return DocumentSpec(
        title="Otters SC Season Report",
        subtitle="2025/26",
        kind="document",
        doc_format="season_report",
        source_refs=["run:r1", "meet results file"],
        sections=[
            Section(
                blocks=[
                    m.heading("Highlights", 1),
                    m.text("A **strong** season."),
                    m.bullet_list(["Grew to 120 members", "9 medals"]),
                ]
            ),
            Section(
                blocks=[
                    m.heading("By the numbers", 2),
                    m.kpi_row([{"value": "42", "label": "PBs"}, {"value": "9", "label": "Medals"}]),
                    m.table(
                        ["Swimmer", "PBs"], [["Ada", "5"], ["Bo", "4"]], caption="Top PB makers"
                    ),
                ]
            ),
        ],
    )


def _deck():
    return DocumentSpec(
        title="AGM 2026",
        kind="deck",
        doc_format="agm_deck",
        sections=[
            Section(layout="cover", blocks=[m.heading("AGM 2026", 1)]),
            Section(blocks=[m.heading("The year", 2), m.bullet_list(["120 members", "9 medals"])]),
            Section(layout="closing", blocks=[m.heading("Thank you", 1)]),
        ],
    )


def test_docx_export_is_valid_and_has_content(tmp_path):
    out = export.document_docx(_doc(), tmp_path / "report.docx")
    assert out.exists()
    import docx

    doc = docx.Document(str(out))
    text = "\n".join(p.text for p in doc.paragraphs)
    assert "Otters SC Season Report" in text
    assert "strong season" in text  # markup stripped to runs
    assert "Grew to 120 members" in text
    # the table came across as a real editable table
    assert len(doc.tables) == 1
    cells = [c.text for row in doc.tables[0].rows for c in row.cells]
    assert "Swimmer" in cells and "Ada" in cells
    # sources footer
    assert "Sources:" in text


def test_pptx_export_one_slide_per_section(tmp_path):
    out = export.document_pptx(_deck(), tmp_path / "agm.pptx")
    assert out.exists()
    from pptx import Presentation

    prs = Presentation(str(out))
    assert len(prs.slides) == 3
    # text made it onto the slides
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text_frame.text)
    joined = "\n".join(all_text)
    assert "AGM 2026" in joined
    assert "120 members" in joined


def test_pptx_export_includes_tables(tmp_path):
    out = export.document_pptx(_doc(), tmp_path / "report.pptx")
    from pptx import Presentation

    prs = Presentation(str(out))
    found_table = any(shape.has_table for slide in prs.slides for shape in slide.shapes)
    assert found_table


def test_export_formats_listed():
    assert export.EXPORT_FORMATS == ("pdf", "pptx", "docx")


# ---------------------------------------------------------------------------
# Security: an export must only embed images that live under DATA_DIR. Specs are
# tenant-editable (the advanced JSON editor), so an absolute path to an arbitrary
# server file — or another tenant's assets — must never be baked into a DOCX/PPTX
# (cross-tenant read / local-file disclosure). Mirrors render._img_src's guard.
# ---------------------------------------------------------------------------


def _png(path, colour):
    from PIL import Image

    Image.new("RGB", (48, 32), colour).save(path)
    return path


def _media_doc(*srcs):
    return DocumentSpec(
        title="pics",
        sections=[Section(blocks=[m.media(str(s)) for s in srcs])],
    )


def _docx_media_bytes(out):
    import zipfile

    with zipfile.ZipFile(out) as z:
        return {z.read(n) for n in z.namelist() if n.startswith("word/media/")}


def _pptx_media_bytes(out):
    import zipfile

    with zipfile.ZipFile(out) as z:
        return {z.read(n) for n in z.namelist() if n.startswith("ppt/media/")}


def test_docx_export_drops_image_outside_data_dir(tmp_path, monkeypatch):
    data_dir = tmp_path / "data"
    data_dir.mkdir()
    monkeypatch.setenv("DATA_DIR", str(data_dir))
    outside = _png(tmp_path / "secret_outside.png", "magenta")  # sibling of DATA_DIR
    inside = _png(data_dir / "legit.png", "green")

    out = export.document_docx(_media_doc(outside, inside), tmp_path / "o.docx")
    embedded = _docx_media_bytes(out)
    assert outside.read_bytes() not in embedded  # cross-tenant/LFI src refused
    assert inside.read_bytes() in embedded  # legitimate DATA_DIR image still embeds


def test_pptx_export_drops_image_outside_data_dir(tmp_path, monkeypatch):
    data_dir = tmp_path / "data"
    data_dir.mkdir()
    monkeypatch.setenv("DATA_DIR", str(data_dir))
    outside = _png(tmp_path / "secret_outside.png", "magenta")
    inside = _png(data_dir / "legit.png", "green")

    out = export.document_pptx(_media_doc(outside, inside), tmp_path / "o.pptx")
    embedded = _pptx_media_bytes(out)
    assert outside.read_bytes() not in embedded
    assert inside.read_bytes() in embedded


def test_export_refuses_remote_and_traversal_srcs(tmp_path, monkeypatch):
    data_dir = tmp_path / "data"
    data_dir.mkdir()
    monkeypatch.setenv("DATA_DIR", str(data_dir))
    # http(s) (no SSRF), data: URIs, and ../ traversal all resolve to nothing.
    for bad in (
        "https://evil.example/x.png",
        "http://169.254.169.254/latest",
        "data:image/png;base64,AAAA",
        "../secret_outside.png",
        str(tmp_path / ".." / "etc" / "hostname"),
    ):
        assert export._img_path(m.media(bad)) is None
