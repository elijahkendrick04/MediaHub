"""documents.import_doc — bring an existing file in as editable blocks (roadmap 1.15).

"Edit PDF / import PowerPoint / import Word", done honestly: a PPTX, DOCX or PDF is
parsed into a :class:`~documents.models.DocumentSpec` of text / heading / list /
table blocks you can then edit and re-export. **Fidelity is bounded and stated** —
we recover the structured text and tables (the parts worth editing), not a
pixel-perfect re-creation of the original layout, fonts or vector art. Images are
not round-tripped (noted on import). When a parser dependency is missing the caller
gets a clear error.
"""

from __future__ import annotations

import re
import zipfile
from pathlib import Path

from . import models as m
from .models import DocumentSpec, Section

# The level of fidelity callers should expect, surfaced in the UI.
FIDELITY_NOTE = (
    "Imported as editable text and tables (bounded fidelity): original layout, "
    "fonts, images and vector art are not preserved."
)

# Defensive caps for untrusted uploads (CLAUDE.md security focus: zip bombs).
# docx/pptx are ZIP archives: only the 50 MB request cap bounds the *compressed*
# input, so a tiny archive could still declare gigabytes uncompressed and OOM
# the worker inside python-docx/pptx. These limits mirror the intent of
# ``interpreter._zip_safety`` with caps generous enough for real office files
# (a big deck holds thousands of members and tens of MB of XML/media).
_MAX_ZIP_MEMBERS = 10_000
_MAX_MEMBER_BYTES = 256 * 1024 * 1024
_MAX_TOTAL_BYTES = 512 * 1024 * 1024
_MAX_MEMBER_RATIO = 200  # checked only for members > 1 MB uncompressed
_MAX_PDF_PAGES = 500


def _check_office_zip(path: str | Path, kind: str) -> None:
    """Reject a docx/pptx whose central directory declares a decompression bomb.

    Inspects declared sizes only — no member bytes are decompressed. Raises
    ``ValueError`` (the import route surfaces it as an honest 422)."""
    try:
        with zipfile.ZipFile(str(path)) as zf:
            infos = zf.infolist()
    except zipfile.BadZipFile as e:
        raise ValueError(f"not a valid {kind} file: {e}") from e
    if len(infos) > _MAX_ZIP_MEMBERS:
        raise ValueError(
            f"{kind} archive has {len(infos)} members (limit {_MAX_ZIP_MEMBERS}) — refusing import"
        )
    total = 0
    for info in infos:
        if info.is_dir():
            continue
        usize = int(info.file_size)
        csize = max(1, int(info.compress_size or 1))
        if usize > _MAX_MEMBER_BYTES:
            raise ValueError(
                f"{kind} member {info.filename!r} declares {usize} uncompressed bytes "
                f"(limit {_MAX_MEMBER_BYTES}) — refusing import"
            )
        if usize > 1024 * 1024 and usize // csize > _MAX_MEMBER_RATIO:
            raise ValueError(
                f"{kind} member {info.filename!r} has compression ratio "
                f"{usize // csize}:1 (limit {_MAX_MEMBER_RATIO}:1) — refusing import"
            )
        total += usize
        if total > _MAX_TOTAL_BYTES:
            raise ValueError(
                f"{kind} archive declares over {_MAX_TOTAL_BYTES} total uncompressed bytes "
                "— refusing import"
            )


def _heading_level(style_name: str) -> int:
    mtc = re.search(r"(\d+)", style_name or "")
    if mtc:
        return max(1, min(3, int(mtc.group(1))))
    return 1


def import_docx(path: str | Path) -> DocumentSpec:
    """Parse a Word .docx into a flowing document of heading/text/list/table blocks."""
    try:
        import docx
        from docx.table import Table
        from docx.text.paragraph import Paragraph
    except Exception as e:  # pragma: no cover
        raise RuntimeError(f"DOCX import needs python-docx: {e}") from e

    _check_office_zip(path, "DOCX")
    doc = docx.Document(str(path))
    blocks: list[m.Block] = []
    body = doc.element.body
    for child in body.iterchildren():
        tag = child.tag.split("}")[-1]
        if tag == "p":
            para = Paragraph(child, doc)
            txt = para.text.strip()
            if not txt:
                continue
            style = (para.style.name if para.style else "") or ""
            if style.startswith("Heading"):
                blocks.append(m.heading(txt, _heading_level(style)))
            elif style.startswith("List"):
                blocks.append(m.bullet_list([txt], ordered="Number" in style))
            else:
                blocks.append(m.text(txt))
        elif tag == "tbl":
            table = Table(child, doc)
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            if not rows:
                continue
            blocks.append(m.table(rows[0], rows[1:]))

    title = (
        (doc.core_properties.title or Path(path).stem) if doc.core_properties else Path(path).stem
    )
    return DocumentSpec(
        title=str(title) or "Imported document",
        kind="document",
        doc_format="blank",
        geometry="a4",
        meta={"imported_from": "docx", "fidelity": FIDELITY_NOTE},
        sections=[Section(blocks=blocks or [m.text("")])],
    )


def import_pptx(path: str | Path) -> DocumentSpec:
    """Parse a PowerPoint .pptx into a deck — one section per slide (text + tables)."""
    try:
        from pptx import Presentation
    except Exception as e:  # pragma: no cover
        raise RuntimeError(f"PPTX import needs python-pptx: {e}") from e

    _check_office_zip(path, "PPTX")
    prs = Presentation(str(path))
    sections: list[Section] = []
    for slide in prs.slides:
        blocks: list[m.Block] = []
        title_text = ""
        try:
            if slide.shapes.title and slide.shapes.title.text.strip():
                title_text = slide.shapes.title.text.strip()
                blocks.append(m.heading(title_text, 2))
        except (AttributeError, ValueError):
            pass
        for shape in slide.shapes:
            if shape.has_table:
                tbl = shape.table
                rows = [[cell.text.strip() for cell in row.cells] for row in tbl.rows]
                if rows:
                    blocks.append(m.table(rows[0], rows[1:]))
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                txt = "".join(run.text for run in para.runs).strip() or para.text.strip()
                if not txt or txt == title_text:
                    continue
                if para.level and para.level > 0:
                    blocks.append(m.bullet_list([txt]))
                else:
                    blocks.append(m.text(txt))
        sections.append(Section(blocks=blocks or [m.text("")]))

    return DocumentSpec(
        title=Path(path).stem or "Imported deck",
        kind="deck",
        doc_format="blank",
        geometry="slide_16_9",
        meta={"imported_from": "pptx", "fidelity": FIDELITY_NOTE},
        sections=sections or [Section(blocks=[m.text("")])],
    )


def import_pdf(path: str | Path) -> DocumentSpec:
    """Parse a PDF into a flowing document — one section per page (text + tables)."""
    try:
        import pdfplumber
    except Exception as e:  # pragma: no cover
        raise RuntimeError(f"PDF import needs pdfplumber: {e}") from e

    sections: list[Section] = []
    with pdfplumber.open(str(path)) as pdf:
        if len(pdf.pages) > _MAX_PDF_PAGES:
            raise ValueError(
                f"PDF has {len(pdf.pages)} pages (limit {_MAX_PDF_PAGES}) — refusing import"
            )
        for page in pdf.pages:
            blocks: list[m.Block] = []
            try:
                for tbl in page.extract_tables() or []:
                    clean = [[("" if c is None else str(c)).strip() for c in row] for row in tbl]
                    if clean:
                        blocks.append(m.table(clean[0], clean[1:]))
            except Exception:
                pass
            text = (page.extract_text() or "").strip()
            if text:
                # split into paragraphs on blank lines, keep it simple
                for chunk in re.split(r"\n\s*\n", text):
                    chunk = chunk.strip()
                    if chunk:
                        blocks.append(m.text(chunk))
            sections.append(Section(blocks=blocks or [m.text("")], break_before=bool(sections)))

    return DocumentSpec(
        title=Path(path).stem or "Imported PDF",
        kind="document",
        doc_format="blank",
        geometry="a4",
        meta={"imported_from": "pdf", "fidelity": FIDELITY_NOTE},
        sections=sections or [Section(blocks=[m.text("")])],
    )


def import_file(path: str | Path) -> DocumentSpec:
    """Dispatch by extension to the right importer."""
    ext = Path(path).suffix.lower().lstrip(".")
    if ext == "docx":
        return import_docx(path)
    if ext == "pptx":
        return import_pptx(path)
    if ext == "pdf":
        return import_pdf(path)
    raise ValueError(f"unsupported import type: .{ext} (supported: pdf, docx, pptx)")


__all__ = ["FIDELITY_NOTE", "import_file", "import_docx", "import_pptx", "import_pdf"]
