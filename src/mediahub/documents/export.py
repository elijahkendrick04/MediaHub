"""documents.export — take-it-elsewhere editing exports (roadmap 1.15).

A club sometimes needs to hand a document to someone who wants to *edit* it in
PowerPoint or Word. These deterministic converters (python-pptx / python-docx)
emit a faithful-but-bounded editable file from a :class:`~documents.models.DocumentSpec`:
the text, headings, lists, tables and KPI lines come across as native editable
objects; charts and photos are placed as images. Fidelity is **bounded and stated**
— this is "open it elsewhere and tweak it", not a pixel-perfect re-layout.

The canonical, on-brand output is always the PDF (:func:`documents.render`); these
are the convenience exits. When python-pptx/docx isn't installed the caller gets a
clear error (the PDF path still works).
"""

from __future__ import annotations

import io
from pathlib import Path
from typing import Any, Optional, Union

from .models import Block, DocumentSpec

# Cap on a decoded ``data:`` image before embedding — a hand-authored spec could
# carry a huge base64 blob, and python-docx/pptx would hold the whole thing in
# memory. Generous for a real photo, bounded against an OOM (mirrors the intent
# of import_doc's zip-bomb caps).
_MAX_EMBED_BYTES = 25 * 1024 * 1024

# An embeddable image source for python-docx/pptx: a filesystem path (str) or an
# in-memory stream (a decoded ``data:`` image).
ImageSource = Union[str, io.BytesIO]

# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------


def _text_lines(block: Block) -> list[str]:
    """Flatten a text-like block to plain lines (markup stripped)."""
    import re

    p = block.props or {}
    if block.kind == "heading":
        return [str(p.get("text", ""))]
    if block.kind == "text":
        raw = str(p.get("text", ""))
        return [re.sub(r"\*+", "", raw)]
    if block.kind == "list":
        return [f"• {i}" for i in (p.get("items") or [])]
    if block.kind == "quote":
        who = p.get("attribution")
        line = f"“{p.get('text', '')}”"
        return [line + (f" - {who}" if who else "")]
    if block.kind == "stat":
        return [f"{p.get('value', '')} - {p.get('label', '')}"]
    if block.kind == "kpi_row":
        return [f"{s.get('value', '')} - {s.get('label', '')}" for s in (p.get("stats") or [])]
    return []


def _chart_png(block: Block, brand_kit: Any, role_vars: Optional[dict]) -> Optional[Path]:
    """Render a chart block to a PNG for embedding (None on any failure)."""
    try:
        from mediahub.charts.export import chart_png_path
        from mediahub.charts.models import ChartSpec

        spec = ChartSpec.from_dict((block.props or {}).get("chart") or {})
        if spec is None:
            return None
        return chart_png_path(spec, fmt="landscape", role_vars=role_vars, brand_kit=brand_kit)
    except Exception:
        return None  # bounded fidelity: skip the image, keep the document


def _img_source(block: Block) -> Optional[ImageSource]:
    """Resolve a block image src to something python-docx/pptx can embed.

    - A ``data:image/...;base64,...`` URI is decoded to an in-memory stream
      (size-capped), so a hand-authored or imported image survives the export —
      matching the render path, which also embeds ``data:`` images.
    - A file path is honoured only when it resolves *inside* ``DATA_DIR``. Specs
      are tenant-editable (the advanced JSON editor), so an absolute path to an
      arbitrary server file — or another tenant's assets under
      ``DATA_DIR/<other>`` — must never be baked into a DOCX/PPTX (cross-tenant
      read / local file disclosure); this mirrors ``render._img_src``.
    - Remote URLs are never fetched (no server-side SSRF).

    Returns a str path or a ``BytesIO``, or ``None`` to skip the image."""
    import base64
    import binascii
    import os

    src = str((block.props or {}).get("src", "")).strip()
    if not src:
        return None
    low = src.lower()
    if low.startswith("data:image/"):
        header, _, payload = src.partition(",")
        if "base64" not in header.lower() or not payload:
            return None  # only base64 image payloads are embeddable here
        try:
            raw = base64.b64decode(payload, validate=False)
        except (ValueError, binascii.Error):
            return None
        if not raw or len(raw) > _MAX_EMBED_BYTES:
            return None
        return io.BytesIO(raw)
    if low.startswith(("http://", "https://", "data:")):
        return None  # remote (SSRF) or a non-image data: URI — skip, never fetch
    if low.startswith("file:"):
        from urllib.parse import unquote, urlparse

        try:
            src = unquote(urlparse(src).path)
        except (OSError, ValueError):
            return None
    try:
        rp = Path(src).resolve()
        root = Path(os.environ.get("DATA_DIR", ".")).resolve()
        if rp.is_file() and rp.is_relative_to(root):
            return str(rp)
    except (OSError, ValueError):
        return None
    return None


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------


def document_docx(
    spec: DocumentSpec,
    out_path: str | Path,
    *,
    brand_kit: Any = None,
    role_vars: Optional[dict[str, str]] = None,
) -> Path:
    """Export a DocumentSpec to an editable Word (.docx) file (bounded fidelity)."""
    try:
        import docx
        from docx.shared import Inches, Pt
    except Exception as e:  # pragma: no cover
        raise RuntimeError(f"DOCX export needs python-docx: {e}") from e

    doc = docx.Document()
    doc.core_properties.title = spec.title
    doc.add_heading(spec.title, level=0)
    if spec.subtitle:
        doc.add_paragraph(spec.subtitle)

    for section in spec.sections:
        for block in section.blocks:
            _docx_block(doc, block, brand_kit, role_vars, Inches, Pt)

    if spec.source_refs:
        p = doc.add_paragraph()
        run = p.add_run("Sources: " + "; ".join(str(r) for r in spec.source_refs[:8]))
        run.italic = True
        run.font.size = Pt(8)

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out_path))
    return out_path


def _docx_runs(paragraph, raw: str) -> None:
    """Add text to a paragraph, honouring **bold** / *italic* inline markup."""
    import re

    for part in re.split(r"(\*\*[^*]+\*\*|\*[^*]+\*)", raw):
        if not part:
            continue
        if part.startswith("**") and part.endswith("**"):
            paragraph.add_run(part[2:-2]).bold = True
        elif part.startswith("*") and part.endswith("*"):
            paragraph.add_run(part[1:-1]).italic = True
        else:
            paragraph.add_run(part)


def _docx_block(doc, block: Block, brand_kit, role_vars, Inches, Pt) -> None:
    p = block.props or {}
    kind = block.kind
    if kind == "heading":
        doc.add_heading(str(p.get("text", "")), level=max(1, min(3, int(p.get("level", 2) or 2))))
    elif kind == "text":
        _docx_runs(doc.add_paragraph(), str(p.get("text", "")))
    elif kind == "list":
        style = "List Number" if p.get("ordered") else "List Bullet"
        for item in p.get("items") or []:
            try:
                doc.add_paragraph(str(item), style=style)
            except KeyError:
                doc.add_paragraph(f"• {item}")
    elif kind == "quote":
        para = doc.add_paragraph()
        run = para.add_run(f"“{p.get('text', '')}”")
        run.italic = True
        if p.get("attribution"):
            doc.add_paragraph(f"- {p['attribution']}")
    elif kind in ("stat", "kpi_row"):
        for line in _text_lines(block):
            doc.add_paragraph(line)
    elif kind == "table":
        _docx_table(doc, p, Pt)
    elif kind == "chart":
        png = _chart_png(block, brand_kit, role_vars)
        if png:
            doc.add_picture(str(png), width=Inches(6.0))
    elif kind in ("card", "media"):
        img = _img_source(block)
        if img is not None:
            try:
                doc.add_picture(img, width=Inches(5.0))
            except Exception:
                pass  # bounded fidelity: skip an unembeddable image, keep the doc
        if p.get("caption"):
            doc.add_paragraph(str(p["caption"]))
    elif kind == "columns":
        for col in p.get("columns") or []:
            for b in col or []:
                _docx_block(doc, Block.from_dict(b), brand_kit, role_vars, Inches, Pt)


def _docx_table(doc, p: dict, Pt) -> None:
    cols = p.get("columns") or []
    rows = p.get("rows") or []
    if not cols and not rows:
        return
    ncols = len(cols) or (len(rows[0]) if rows else 1)
    table = doc.add_table(rows=0, cols=ncols)
    try:
        table.style = "Light Grid Accent 1"
    except KeyError:
        pass
    if cols:
        hdr = table.add_row().cells
        for i, c in enumerate(cols[:ncols]):
            hdr[i].text = str(c)
    for r in rows:
        cells = table.add_row().cells
        for i, c in enumerate(r[:ncols]):
            cells[i].text = str(c)
    if p.get("caption"):
        cap = doc.add_paragraph(str(p["caption"]))
        cap.runs[0].italic = True
        cap.runs[0].font.size = Pt(8)


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------


def document_pptx(
    spec: DocumentSpec,
    out_path: str | Path,
    *,
    brand_kit: Any = None,
    role_vars: Optional[dict[str, str]] = None,
) -> Path:
    """Export a DocumentSpec to an editable PowerPoint (.pptx) — one slide per
    section (bounded fidelity: a simple vertical stack, fully editable)."""
    try:
        from pptx import Presentation
        from pptx.util import Emu, Inches, Pt
    except Exception as e:  # pragma: no cover
        raise RuntimeError(f"PPTX export needs python-pptx: {e}") from e

    prs = Presentation()
    # Slide size: 16:9 for decks, A4-ish portrait for documents.
    if spec.is_deck and spec.geometry == "slide_4_3":
        prs.slide_width, prs.slide_height = Inches(10), Inches(7.5)
    elif spec.is_deck:
        prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    else:
        prs.slide_width, prs.slide_height = Inches(8.27), Inches(11.69)  # A4

    blank = prs.slide_layouts[6]  # fully blank
    margin = Inches(0.5)
    content_w = prs.slide_width - 2 * margin

    for section in spec.sections:
        slide = prs.slides.add_slide(blank)
        y = margin
        # Collect text-like content into one frame, place images/tables below.
        text_lines: list[tuple[str, bool]] = []  # (line, is_heading)
        media: list[Block] = []
        tables: list[dict] = []
        for block in section.blocks:
            if block.kind in ("heading", "text", "list", "quote", "stat", "kpi_row"):
                is_head = block.kind == "heading"
                for line in _text_lines(block):
                    text_lines.append((line, is_head))
            elif block.kind == "table":
                tables.append(block.props or {})
            elif block.kind in ("chart", "card", "media"):
                media.append(block)
            elif block.kind == "columns":
                for col in (block.props or {}).get("columns") or []:
                    for b in col or []:
                        bb = Block.from_dict(b)
                        if bb.kind in ("chart", "card", "media"):
                            media.append(bb)
                        else:
                            for line in _text_lines(bb):
                                text_lines.append((line, bb.kind == "heading"))

        if text_lines:
            box = slide.shapes.add_textbox(margin, y, content_w, Inches(2.2))
            tf = box.text_frame
            tf.word_wrap = True
            first = True
            for line, is_head in text_lines:
                para = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                para.text = line
                para.font.size = Pt(28) if is_head else Pt(16)
                para.font.bold = is_head
            y = y + Inches(2.4)

        for tprops in tables:
            y = _pptx_table(slide, tprops, margin, y, content_w, Inches, Pt, Emu)

        for block in media:
            if block.kind == "chart":
                cp = _chart_png(block, brand_kit, role_vars)
                src = str(cp) if cp else None
            else:
                src = _img_source(block)
            if src is not None and y < prs.slide_height - Inches(1):
                try:
                    slide.shapes.add_picture(src, margin, y, width=content_w)
                    y = y + Inches(3.0)
                except Exception:
                    pass  # bounded fidelity: skip an unembeddable image, keep the slide

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def _pptx_table(slide, p: dict, x, y, width, Inches, Pt, Emu):
    cols = p.get("columns") or []
    rows = p.get("rows") or []
    ncols = len(cols) or (len(rows[0]) if rows else 1)
    nrows = (1 if cols else 0) + len(rows)
    if nrows == 0 or ncols == 0:
        return y
    height = Inches(0.3 * nrows)
    shape = slide.shapes.add_table(nrows, ncols, x, y, width, height)
    table = shape.table
    r = 0
    if cols:
        for ci, c in enumerate(cols[:ncols]):
            cell = table.cell(0, ci)
            cell.text = str(c)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(11)
                para.font.bold = True
        r = 1
    for row in rows:
        for ci, c in enumerate(row[:ncols]):
            cell = table.cell(r, ci)
            cell.text = str(c)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(11)
        r += 1
    return y + height + Inches(0.2)


EXPORT_FORMATS = ("pdf", "pptx", "docx")


__all__ = ["EXPORT_FORMATS", "document_docx", "document_pptx"]
